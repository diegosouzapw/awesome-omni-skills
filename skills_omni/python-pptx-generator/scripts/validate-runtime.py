#!/usr/bin/env python3
"""Inspect a PowerPoint template for slide layouts and placeholder metadata.

Usage:
    python scripts/validate-runtime.py path/to/template.pptx
"""

from __future__ import annotations

import argparse
import sys
from pathlib import Path

from pptx import Presentation

EMU_PER_INCH = 914400


def emu_to_inches(value: int) -> str:
    return f"{value / EMU_PER_INCH:.2f}in"


def format_placeholder(ph) -> str:
    pf = ph.placeholder_format
    name = getattr(ph, "name", "<unnamed>")
    ph_type = str(pf.type)
    return (
        f"idx={pf.idx} | name={name} | type={ph_type} | "
        f"left={emu_to_inches(ph.left)} top={emu_to_inches(ph.top)} "
        f"width={emu_to_inches(ph.width)} height={emu_to_inches(ph.height)}"
    )


def inspect_template(path: Path) -> int:
    if not path.exists() or not path.is_file():
        print(f"ERROR: file not found: {path}", file=sys.stderr)
        return 2
    if path.suffix.lower() != ".pptx":
        print(f"ERROR: expected a .pptx file: {path}", file=sys.stderr)
        return 2

    try:
        prs = Presentation(str(path))
    except Exception as exc:  # pragma: no cover - defensive runtime path
        print(f"ERROR: could not open presentation: {exc}", file=sys.stderr)
        return 3

    print(f"Template: {path}")
    print(f"Slide size: {emu_to_inches(prs.slide_width)} x {emu_to_inches(prs.slide_height)}")
    print(f"Layout count: {len(prs.slide_layouts)}")
    print()

    for i, layout in enumerate(prs.slide_layouts):
        placeholders = list(layout.placeholders)
        print(f"[{i}] Layout: {layout.name}")
        if not placeholders:
            print("  (no placeholders)")
            print()
            continue

        for ph in placeholders:
            try:
                line = format_placeholder(ph)
            except Exception as exc:  # pragma: no cover - defensive runtime path
                line = f"<placeholder inspection failed: {exc}>"
            print(f"  - {line}")
        print()

    return 0


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(
        description="Inspect a .pptx template and print layout and placeholder details."
    )
    parser.add_argument("pptx_path", help="Path to the template .pptx file")
    return parser.parse_args()


def main() -> int:
    args = parse_args()
    return inspect_template(Path(args.pptx_path))


if __name__ == "__main__":
    sys.exit(main())
